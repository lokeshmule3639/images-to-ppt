import os
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import imagehash

app = Flask(__name__)
CORS(app)  # Allow frontend to access backend

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def is_duplicate(image_path, seen_hashes):
    """Check if an image is a duplicate using perceptual hashing."""
    img_hash = imagehash.average_hash(Image.open(image_path))
    if img_hash in seen_hashes:
        return True
    seen_hashes.add(img_hash)
    return False

@app.route("/upload", methods=["POST"])
def upload_files():
    """Handle file uploads and generate a PowerPoint presentation."""
    files = request.files.getlist("images")
    seen_hashes = set()
    
    ppt = Presentation()
    ppt_file = os.path.join(UPLOAD_FOLDER, "WhatsApp_Images_Presentation.pptx")

    positions = [(Inches(0.5), Inches(0.5)), (Inches(5.0), Inches(0.5)), 
                 (Inches(0.5), Inches(4.0)), (Inches(5.0), Inches(4.0))]
    img_width, img_height = Inches(4.0), Inches(3.0)

    slide = None
    count = 0

    for file in files:
        file_path = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(file_path)

        if is_duplicate(file_path, seen_hashes):
            os.remove(file_path)  # Remove duplicate
            continue

        if count % 4 == 0:  # New slide after every 4 images
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
        left, top = positions[count % 4]
        slide.shapes.add_picture(file_path, left, top, width=img_width, height=img_height)

        count += 1

    ppt.save(ppt_file)
    return send_file(ppt_file, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
